import os
import pickle
import base64
from email.mime.text import MIMEText
import streamlit as st
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
import google.generativeai as genai
from langgraph.graph import StateGraph, END
from typing import TypedDict
from google.oauth2 import service_account
import smtplib
from dotenv import load_dotenv
import re
import io
from googleapiclient.http import MediaIoBaseDownload
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
import csv

# --- Streamlit Page Config ---
st.set_page_config(page_title="Google Workspace Agent App", page_icon="🤖")
st.title("Google Workspace Agent App")
st.markdown("""
Summarize a Google Doc, Sheet, Drive, PDF, DOCX, XLSX, PPTX, or CSV and email the summary using an agent-based workflow!\
**Your data is processed securely.**
""")

st.info("""
**Instructions:**
1. Select the Google file type (Doc, Sheet, Drive, PDF, DOCX, XLSX, PPTX, or CSV).
2. Paste the shareable link (ensure the file is accessible to the app's Google account).
3. Click 'Summarize and Email (Agent)'.
""")

# --- Load secrets from Streamlit secrets ---
GEMINI_API_KEY = st.secrets.get("GEMINI_API_KEY")
TARGET_EMAIL = st.secrets["TARGET_EMAIL"]
GOOGLE_CREDENTIALS = st.secrets["GOOGLE_CREDENTIALS"]  # Should be the JSON string

if not GEMINI_API_KEY:
    # Fallback: try to load from .env (for local development)
    load_dotenv()
    GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

if not GEMINI_API_KEY:
    st.error("GEMINI_API_KEY not found in Streamlit secrets or .env file.")
    st.stop()

SCOPES = [
    'https://www.googleapis.com/auth/drive.readonly',
    'https://www.googleapis.com/auth/gmail.send',
    'https://www.googleapis.com/auth/spreadsheets.readonly',
    'https://www.googleapis.com/auth/presentations.readonly',
    'https://www.googleapis.com/auth/documents.readonly'
]

# --- Write Google credentials to file at runtime ---
CREDENTIALS_PATH = "credentials.json"
with open(CREDENTIALS_PATH, "w") as f:
    f.write(GOOGLE_CREDENTIALS)

# --- Google Auth ---
def get_credentials():
    credentials = service_account.Credentials.from_service_account_file(
        "credentials.json",
        scopes=[
            'https://www.googleapis.com/auth/drive.readonly',
            'https://www.googleapis.com/auth/gmail.send',
            'https://www.googleapis.com/auth/spreadsheets.readonly',
            'https://www.googleapis.com/auth/presentations.readonly',
            'https://www.googleapis.com/auth/documents.readonly'
        ]
    )
    return credentials

# --- Extraction Functions ---
def extract_doc_text(docs_service, doc_id):
    try:
        doc = docs_service.documents().get(documentId=doc_id).execute()
        text = []
        for element in doc.get('body', {}).get('content', []):
            if 'paragraph' in element:
                for p_elem in element['paragraph'].get('elements', []):
                    if 'textRun' in p_elem:
                        text.append(p_elem['textRun'].get('content', ''))
        return ''.join(text)
    except Exception as e:
        return f"Error reading Google Doc: {e}"

def extract_sheet_text(sheets_service, sheet_id):
    try:
        result = sheets_service.spreadsheets().values().get(
            spreadsheetId=sheet_id, range='A1:Z1000').execute()
        values = result.get('values', [])
        sheet_data = ""
        if values:
            for row in values:
                sheet_data += "\t".join([str(cell) for cell in row]) + "\n"
        return sheet_data
    except Exception as e:
        return f"Error reading Google Sheet: {e}"

def extract_presentation_text(slides_service, presentation_id):
    try:
        presentation = slides_service.presentations().get(presentationId=presentation_id).execute()
        slides = presentation.get('slides', [])
        text_content = []
        for i, slide in enumerate(slides):
            elements = slide.get('pageElements', [])
            for element in elements:
                shape = element.get('shape')
                if shape:
                    text_elements = shape.get('text', {}).get('textElements', [])
                    for te in text_elements:
                        if 'textRun' in te:
                            text_content.append(te['textRun'].get('content', ''))
        return '\n'.join(text_content)
    except Exception as e:
        return f"Error reading Google Slides: {e}"

def extract_pdf_text(drive_service, file_id):
    try:
        request = drive_service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
        fh.seek(0)
        reader = PyPDF2.PdfReader(fh)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        return f"Error reading PDF: {e}"

# Extraction functions for Office files

def extract_docx_text(drive_service, file_id):
    try:
        request = drive_service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
        fh.seek(0)
        doc = docx.Document(fh)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        return f"Error reading DOCX: {e}"

def extract_xlsx_text(drive_service, file_id):
    try:
        request = drive_service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
        fh.seek(0)
        wb = openpyxl.load_workbook(fh, data_only=True)
        text = ""
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text += "\t".join([str(cell) if cell is not None else '' for cell in row]) + "\n"
        return text
    except Exception as e:
        return f"Error reading XLSX: {e}"

def extract_pptx_text(drive_service, file_id):
    try:
        request = drive_service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
        fh.seek(0)
        prs = Presentation(fh)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error reading PPTX: {e}"

def extract_csv_text(drive_service, file_id):
    try:
        request = drive_service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
        fh.seek(0)
        text = ""
        fh.seek(0)
        decoded = fh.read().decode('utf-8', errors='replace').splitlines()
        reader = csv.reader(decoded)
        for row in reader:
            text += ", ".join(row) + "\n"
        return text
    except Exception as e:
        return f"Error reading CSV: {e}"

# --- LangGraph State Schema ---
class AgentState(TypedDict):
    file_type: str
    file_id: str
    extracted_text: str
    summary: str
    email_status: str
    recipient_email: str

# --- LangGraph Nodes ---
def fetch_content_node(state: AgentState):
    creds = get_credentials()
    file_type = state['file_type']
    file_id = state['file_id']
    text = ""
    if file_type == "Doc":
        docs_service = build('docs', 'v1', credentials=creds)
        text = extract_doc_text(docs_service, file_id)
    elif file_type == "Sheet":
        sheets_service = build('sheets', 'v4', credentials=creds)
        text = extract_sheet_text(sheets_service, file_id)
    elif file_type == "Slides":
        slides_service = build('slides', 'v1', credentials=creds)
        text = extract_presentation_text(slides_service, file_id)
    elif file_type == "PDF":
        drive_service = build('drive', 'v3', credentials=creds)
        text = extract_pdf_text(drive_service, file_id)
    elif file_type == "DOCX":
        drive_service = build('drive', 'v3', credentials=creds)
        text = extract_docx_text(drive_service, file_id)
    elif file_type == "XLSX":
        drive_service = build('drive', 'v3', credentials=creds)
        text = extract_xlsx_text(drive_service, file_id)
    elif file_type == "PPTX":
        drive_service = build('drive', 'v3', credentials=creds)
        text = extract_pptx_text(drive_service, file_id)
    elif file_type == "CSV":
        drive_service = build('drive', 'v3', credentials=creds)
        text = extract_csv_text(drive_service, file_id)
    state['extracted_text'] = text
    return state

def summarize_node(state: AgentState):
    genai.configure(api_key=GEMINI_API_KEY)
    text = state.get('extracted_text', '')
    file_type = state.get('file_type', '')
    # Limit the amount of data sent to Gemini
    max_chars = 4000
    if len(text) > max_chars:
        text = text[:max_chars] + '\n... [truncated]'
    # Tailored prompt for each file type
    if file_type in ["Doc", "DOCX"]:
        prompt = f"Summarize the following document. Focus on the main points, key takeaways, and any action items.\n\nDocument Content:\n{text}\n\nSummary:"
    elif file_type in ["Sheet", "XLSX", "CSV"]:
        prompt = f"The following is tabular data. Provide a summary of the key columns, trends, and any notable data points.\n\nTable Content:\n{text}\n\nSummary:"
    elif file_type in ["Slides", "PPTX"]:
        prompt = f"Summarize the following presentation. Highlight the main topics covered in the slides.\n\nPresentation Content:\n{text}\n\nSummary:"
    elif file_type == "PDF":
        prompt = f"Summarize the following PDF document. Focus on the main ideas and important details.\n\nPDF Content:\n{text}\n\nSummary:"
    else:
        prompt = f"Summarize the following content:\n\n{text}\n\nSummary:"
    try:
        model = genai.GenerativeModel('models/gemini-1.5-pro-latest')
        response = model.generate_content(prompt)
        summary = response.text
    except Exception as e:
        summary = f"Error generating summary: {e}"
    state['summary'] = summary
    return state

def send_email_smtp(subject, body, to_email, smtp_server, smtp_port, smtp_user, smtp_password):
    msg = MIMEText(body)
    msg['Subject'] = subject
    msg['From'] = smtp_user
    msg['To'] = to_email
    try:
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()
            server.login(smtp_user, smtp_password)
            server.sendmail(smtp_user, [to_email], msg.as_string())
        return f"Email sent successfully to {to_email}!"
    except Exception as e:
        return f"Error sending email: {e}"

# --- Streamlit UI ---
file_type = st.selectbox("Select the file type to summarize or list:", ["Doc", "Sheet", "Slides", "Drive", "PDF", "DOCX", "XLSX", "PPTX", "CSV"])
link = st.text_input(f"Paste the Google {file_type} link here:")
recipient_email = st.text_input("Enter the recipient email address:")

st.markdown("---")
use_custom_smtp = st.checkbox("Advanced: Use your own SMTP credentials")

if use_custom_smtp:
    smtp_server = st.text_input("SMTP server", value="smtp.gmail.com")
    smtp_port = st.number_input("SMTP port", value=587)
    smtp_user = st.text_input("SMTP username (your email)")
    smtp_password = st.text_input("SMTP password or app password", type="password")
else:
    smtp_server = st.secrets.get("SMTP_SERVER", "smtp.gmail.com")
    smtp_port = int(st.secrets.get("SMTP_PORT", 587))
    smtp_user = st.secrets.get("SMTP_USER", "")
    smtp_password = st.secrets.get("SMTP_PASSWORD", "")

# Helper to extract ID from link
def extract_id_from_link(link):
    # Handles /d/<id>, id=<id>, and /folders/<id> patterns
    patterns = [
        r"/d/([\w-]+)",
        r"id=([\w-]+)",
        r"/folders/([\w-]+)"
    ]
    for pattern in patterns:
        match = re.search(pattern, link)
        if match:
            return match.group(1)
    return ""

def list_drive_files_recursive(drive_service, folder_id, parent_path=""):
    files = []
    try:
        results = drive_service.files().list(
            q=f"'{folder_id}' in parents and trashed = false",
            fields="files(id, name, mimeType)"
        ).execute()
        for f in results.get('files', []):
            file_path = f"{parent_path}/{f['name']}" if parent_path else f["name"]
            if f['mimeType'] == 'application/vnd.google-apps.folder':
                # Recurse into subfolder
                files.extend(list_drive_files_recursive(drive_service, f['id'], file_path))
            else:
                files.append({"id": f["id"], "name": f["name"], "mimeType": f["mimeType"], "path": file_path})
    except Exception as e:
        st.error(f"Error listing files: {e}")
    return files

if file_type == "Drive" and link:
    folder_id = extract_id_from_link(link)
    if not folder_id:
        st.error("Could not extract folder ID from the link. Please check the link format.")
    else:
        creds = get_credentials()
        drive_service = build('drive', 'v3', credentials=creds)
        files = list_drive_files_recursive(drive_service, folder_id)
        # Supported types for summarization (add Office types)
        supported_types = [
            "application/vnd.google-apps.document",
            "application/vnd.google-apps.spreadsheet",
            "application/vnd.google-apps.presentation",
            "application/pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "text/csv"
        ]
        summaries = {}
        filtered_files = [f for f in files if f["mimeType"] in supported_types]
        col1, col2 = st.columns([1, 2])
        with col1:
            st.markdown("### Files in Folder (and Subfolders)")
            for f in files:
                file_id = f["id"]
                file_name = f["name"]
                file_path = f["path"]
                mime_type = f["mimeType"]
                is_supported = mime_type in supported_types
                st.write(f"{file_path} | Type: {mime_type} {'(Supported)' if is_supported else '(Not Supported)'}")
        # Prepare summaries for supported files
        for f in filtered_files:
            file_id = f['id']
            file_name = f['name']
            file_path = f['path']
            mime_type = f['mimeType']
            summary = None
            if mime_type == "application/vnd.google-apps.document":
                docs_service = build('docs', 'v1', credentials=creds)
                text = extract_doc_text(docs_service, file_id)
                summary = summarize_node({'file_type': 'Doc', 'extracted_text': text, 'file_id': file_id, 'summary': '', 'email_status': '', 'recipient_email': ''})['summary']
            elif mime_type == "application/vnd.google-apps.spreadsheet":
                sheets_service = build('sheets', 'v4', credentials=creds)
                text = extract_sheet_text(sheets_service, file_id)
                summary = summarize_node({'file_type': 'Sheet', 'extracted_text': text, 'file_id': file_id, 'summary': '', 'email_status': '', 'recipient_email': ''})['summary']
            elif mime_type == "application/vnd.google-apps.presentation":
                slides_service = build('slides', 'v1', credentials=creds)
                text = extract_presentation_text(slides_service, file_id)
                summary = summarize_node({'file_type': 'Slides', 'extracted_text': text, 'file_id': file_id, 'summary': '', 'email_status': '', 'recipient_email': ''})['summary']
            elif mime_type == "application/pdf":
                text = extract_pdf_text(drive_service, file_id)
                summary = summarize_node({'file_type': 'PDF', 'extracted_text': text, 'file_id': file_id, 'summary': '', 'email_status': '', 'recipient_email': ''})['summary']
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                text = extract_docx_text(drive_service, file_id)
                summary = summarize_node({'file_type': 'DOCX', 'extracted_text': text, 'file_id': file_id, 'summary': '', 'email_status': '', 'recipient_email': ''})['summary']
            elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                text = extract_xlsx_text(drive_service, file_id)
                summary = summarize_node({'file_type': 'XLSX', 'extracted_text': text, 'file_id': file_id, 'summary': '', 'email_status': '', 'recipient_email': ''})['summary']
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                text = extract_pptx_text(drive_service, file_id)
                summary = summarize_node({'file_type': 'PPTX', 'extracted_text': text, 'file_id': file_id, 'summary': '', 'email_status': '', 'recipient_email': ''})['summary']
            elif mime_type == "text/csv":
                text = extract_csv_text(drive_service, file_id)
                summary = summarize_node({'file_type': 'CSV', 'extracted_text': text, 'file_id': file_id, 'summary': '', 'email_status': '', 'recipient_email': ''})['summary']
            summaries[file_id] = {"title": file_name, "summary": summary, "mime_type": mime_type, "path": file_path}
        with col2:
            st.markdown("### File Summary")
            if filtered_files:
                file_ids = [f["id"] for f in filtered_files]
                file_labels = [summaries[x]["path"] for x in file_ids]
                selected_file_id = st.radio("Select a file to view summary:", file_ids, format_func=lambda x: summaries[x]["path"])
                st.markdown(f"**{summaries[selected_file_id]['title']}**")
                st.write(summaries[selected_file_id]['summary'])
            else:
                st.info("No supported files (Docs, Sheets, Slides, PDF, DOCX, XLSX, PPTX, CSV) found in this folder or subfolders.")
                selected_file_id = None
            col_btn1, col_btn2 = st.columns(2)
            with col_btn1:
                if selected_file_id and st.button("Send summary of selected file"):
                    subject = f"Summary: {summaries[selected_file_id]['title']}"
                    body = f"Title: {summaries[selected_file_id]['title']}\n\nSummary:\n{summaries[selected_file_id]['summary']}"
                    status = send_email_smtp(subject, body, recipient_email, smtp_server, smtp_port, smtp_user, smtp_password)
                    if "successfully" in status.lower():
                        st.success(status)
                    else:
                        st.error(status)
            with col_btn2:
                if summaries and st.button("Send summaries of all files"):
                    subject = "Summaries of all files in Drive folder"
                    body = ""
                    for file_id, info in summaries.items():
                        body += f"Title: {info['title']}\nSummary:\n{info['summary']}\n\n{'-'*40}\n"
                    status = send_email_smtp(subject, body, recipient_email, smtp_server, smtp_port, smtp_user, smtp_password)
                    if "successfully" in status.lower():
                        st.success(status)
                    else:
                        st.error(status)

# Only show the Summarize and Email (Agent) button for non-Drive file types
if file_type != "Drive":
    if st.button("Summarize and Email (Agent)"):
        if not link:
            st.error("Please paste a valid Google link.")
        elif not recipient_email:
            st.error("Please enter a recipient email address.")
        elif use_custom_smtp and (not smtp_user or not smtp_password):
            st.error("Please enter your SMTP username and password.")
        else:
            file_id = extract_id_from_link(link)
            if not file_id:
                st.error("Could not extract file ID from the link. Please check the link format.")
            else:
                actual_file_type = file_type
                if file_type == "Drive":
                    creds = get_credentials()
                    drive_service = build('drive', 'v3', credentials=creds)
                    detected_type = detect_drive_file_type(drive_service, file_id)
                    if not detected_type:
                        st.error("This Drive file is not a supported Google Doc, Sheet, or Slides.")
                        st.stop()
                    actual_file_type = detected_type
                # --- LangGraph Agent Workflow ---
                graph = StateGraph(AgentState)
                graph.add_node('fetch', fetch_content_node)
                graph.add_node('summarize', summarize_node)
                def email_node(state: AgentState):
                    summary = state.get('summary', '')
                    recipient = state.get('recipient_email', '')
                    subject = 'Automated Google Workspace Summary'
                    # Use SMTP credentials from UI or secrets
                    status = send_email_smtp(
                        subject,
                        summary,
                        recipient,
                        smtp_server,
                        smtp_port,
                        smtp_user,
                        smtp_password
                    )
                    state['email_status'] = status
                    return state
                graph.add_node('email', email_node)
                graph.add_edge('fetch', 'summarize')
                graph.add_edge('summarize', 'email')
                graph.add_edge('email', END)
                graph.set_entry_point('fetch')
                app = graph.compile()
                st.info("Running agent workflow...")
                result = app.invoke({
                    "file_type": actual_file_type,
                    "file_id": file_id,
                    "extracted_text": "",
                    "summary": "",
                    "email_status": "",
                    "recipient_email": recipient_email
                })
                st.subheader("Summary:")
                st.write(result.get('summary', ''))
                if "successfully" in result.get('email_status', '').lower():
                    st.success(result.get('email_status', ''))
                else:
                    st.error(result.get('email_status', ''))

SMTP_SERVER = "smtp.gmail.com"
SMTP_PORT = "587"
SMTP_USER = "devarajhansi974@gmail.com"
SMTP_PASSWORD = "pzfgdepfsvtdbfyx"